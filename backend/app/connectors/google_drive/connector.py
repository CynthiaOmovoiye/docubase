"""
Google Drive connector.

Authentication: OAuth access token resolved from the user's ConnectedAccount.

Sync strategy:
  Full sync  — first ingest or when no pageToken is available.
               Lists all text-compatible files in the target folder (or root if
               no folder_id is specified).
  Delta sync — uses the Drive Changes API with a stored pageToken to fetch only
               files changed since the last sync.  result.next_page_token carries
               the new cursor that the ingestion job stores on the Source row.

Supported MIME types: Google Docs (exported as UTF-8 text), Sheets (CSV), Slides (plain),
plain text, Markdown, Office (DOCX/PPTX via binary parse), PDF (binary + pypdf — same as
local ``PDFConnector``).  Single-file ingest also allows ``application/octet-stream`` when the
file name ends in a known extension (Drive often mis-labels PDFs).  Folder sync does not
list generic octet-stream (too noisy).  All ``alt=media`` bodies are decoded from bytes
(never ``response.text`` on PDFs) so chunks stay human-readable.

Security: file_id and folder_id are validated to be Drive IDs (alphanumeric + dash)
before use to prevent accidental traversal.
"""

from __future__ import annotations

import io
import re
from collections.abc import AsyncIterator

import httpx

from app.connectors.base import BaseConnector, ConnectorResult, RawFile
from app.connectors.pdf.text_extract import align_pdf_bytes, extract_readable_pdf_text_from_bytes
from app.core.exceptions import ForbiddenError
from app.core.logging import get_logger

logger = get_logger(__name__)

# Drive resource IDs are alphanumeric + dash + underscore, 20-44 chars
_DRIVE_ID_RE = re.compile(r"^[a-zA-Z0-9_\-]{10,100}$")

# MIME types we can meaningfully ingest
_EXPORTABLE_GOOGLE_TYPES = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}
_NATIVE_TEXT_TYPES = {
    "text/plain",
    "text/markdown",
    "text/x-markdown",
    "application/pdf",  # listed for Drive queries; content fetched via pypdf below
}
# Microsoft Office formats uploaded to Drive — downloaded as binary, parsed locally.
# Maps MIME type → parser label used in _fetch_file_content.
_OFFICE_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
}

# Drive often labels uploads (especially PDFs) as generic binary — still listable / fetchable.
_OCTET_STREAM_MIMES = frozenset(
    {
        "application/octet-stream",
        "binary/octet-stream",
        "application/x-download",
    }
)
_ZIP_MAGIC = b"PK\x03\x04"

_DRIVE_FILES_API = "https://www.googleapis.com/drive/v3/files"
_DRIVE_CHANGES_API = "https://www.googleapis.com/drive/v3/changes"
_DRIVE_START_PAGE_TOKEN_API = "https://www.googleapis.com/drive/v3/changes/startPageToken"

MAX_FILE_SIZE_BYTES = 5_000_000  # 5 MB for Drive files


def _drive_file_extension(name: str) -> str:
    """Lowercase extension including dot (e.g. '.pdf'). Empty if none."""
    base = name.rsplit("/", 1)[-1].split("[", 1)[0].strip().lower()
    if "." not in base:
        return ""
    return "." + base.rsplit(".", 1)[-1]


def _looks_like_binary_mojibake(text: str) -> bool:
    """Reject bodies that are clearly binary mistaken for UTF-8 text."""
    if not text:
        return True
    if text.count("\ufffd") > max(3, len(text) // 150):
        return True
    bad_ctrl = sum(1 for c in text if ord(c) < 32 and c not in "\t\n\r\v\f")
    return len(text) > 400 and bad_ctrl > len(text) // 40


def _decode_drive_file_body(data: bytes, *, name: str, mime: str) -> str | None:
    """
    Normalise a Drive ``alt=media`` download to readable UTF-8 text for chunking.

    Handles: PDF (magic or ``.pdf`` name), DOCX/PPTX (ZIP + python-docx / python-pptx),
    and plain text/markdown with strict rejection of binary mojibake.
    """
    if not data:
        return None
    ext = _drive_file_extension(name)

    def _try_extract_pdf() -> str | None:
        extracted = extract_readable_pdf_text_from_bytes(data, name=name)
        if not extracted or not extracted.strip():
            return None
        if _looks_like_binary_mojibake(extracted):
            logger.warning("gdrive_pdf_extraction_still_binary", name=name)
            return None
        return extracted

    # PDF: never UTF-8-decode raw bytes (mojibake + xref/``%%EOF`` junk in ``chunks``).
    # WPS and similar tools may prefix the file; alignment lives inside the extractor.
    is_declared_pdf = (
        ext == ".pdf"
        or mime == "application/pdf"
        or (mime in _OCTET_STREAM_MIMES and ext == ".pdf")
    )
    if is_declared_pdf:
        return _try_extract_pdf()

    # Drive sometimes serves ``text/plain`` for a real PDF, or a non-``.pdf`` name on PDF bytes.
    # If the body has a standard ``%PDF`` header (possibly after a short prefix) but the name
    # is ``.txt`` / ``.md``, try UTF-8 below when extraction does not produce prose.
    _native_text_exts = {".txt", ".md", ".markdown", ".csv"}
    if align_pdf_bytes(data) is not None and ext not in _native_text_exts:
        return _try_extract_pdf()

    docx_mimes = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }
    pptx_mimes = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }
    if ext in (".docx", ".doc") or mime in docx_mimes:
        if ext in (".docx", ".doc") or (len(data) >= 4 and data[:4] == _ZIP_MAGIC):
            text = _extract_office_text(data, "docx", name)
            if text and text.strip():
                return text
    if ext == ".pptx" or mime in pptx_mimes:
        if ext == ".pptx" or (len(data) >= 4 and data[:4] == _ZIP_MAGIC):
            text = _extract_office_text(data, "pptx", name)
            if text and text.strip():
                return text

    # A PDF mis-labelled as text/plain and named ``.txt`` / ``.md`` still has ``%PDF`` in bytes.
    if align_pdf_bytes(data) is not None and ext in _native_text_exts:
        extracted = _try_extract_pdf()
        if extracted:
            return extracted

    text = data.decode("utf-8-sig", errors="replace").replace("\x00", "")
    text = text.strip()
    if not text:
        return None
    if _looks_like_binary_mojibake(text):
        logger.warning(
            "gdrive_decode_rejected_binary_like",
            name=name,
            mime=mime,
            body_len=len(data),
        )
        return None
    return text


def _validate_drive_id(id_str: str, label: str = "id") -> str:
    if not _DRIVE_ID_RE.match(id_str):
        raise ForbiddenError(f"Invalid Drive {label} {id_str!r}")
    return id_str


def _make_client(token: str) -> httpx.AsyncClient:
    return httpx.AsyncClient(
        headers={"Authorization": f"Bearer {token}"},
        timeout=60.0,
    )


class GoogleDriveConnector(BaseConnector):

    source_type = "google_drive"

    async def validate_connection(
        self, connection_config: dict, access_token: str | None = None
    ) -> bool:
        if not access_token:
            return False
        folder_id = connection_config.get("folder_id")
        if folder_id:
            try:
                _validate_drive_id(folder_id, "folder_id")
            except ForbiddenError:
                return False
        try:
            async with _make_client(access_token) as client:
                resp = await client.get(
                    _DRIVE_FILES_API,
                    params={"pageSize": 1, "fields": "files(id)"},
                )
            return resp.status_code == 200
        except httpx.RequestError:
            return False

    async def fetch(
        self,
        connection_config: dict,
        access_token: str | None = None,
        last_commit_sha: str | None = None,  # unused for Drive
        last_page_token: str | None = None,
    ) -> ConnectorResult:
        if not access_token:
            raise ValueError("Google Drive connector requires an OAuth access token")

        source_id = connection_config.get("source_id", "unknown")
        folder_id = connection_config.get("folder_id")
        file_id = connection_config.get("file_id")

        if folder_id:
            _validate_drive_id(folder_id, "folder_id")
        if file_id:
            _validate_drive_id(file_id, "file_id")

        async with _make_client(access_token) as client:
            # Single-file mode: bypass the folder listing entirely
            if file_id and not folder_id:
                return await self._fetch_single_file(client, source_id, file_id)
            if last_page_token:
                return await self._fetch_delta(client, source_id, last_page_token)
            return await self._fetch_full(client, source_id, folder_id)

    async def stream(
        self,
        connection_config: dict,
        access_token: str | None = None,
    ) -> AsyncIterator[RawFile]:
        result = await self.fetch(connection_config, access_token=access_token)
        for f in result.files:
            yield f

    # ─── Sync implementations ─────────────────────────────────────────────────

    async def _fetch_single_file(
        self,
        client: httpx.AsyncClient,
        source_id: str,
        file_id: str,
    ) -> ConnectorResult:
        """
        Fetch a single Drive file by ID.

        Used when the user picks a specific file (not a folder).
        Skips the folder-listing query entirely — just fetches that one item.
        No delta sync support: single-file sources always do a full refresh.
        """
        logger.info("gdrive_single_file_sync", file_id=file_id)
        errors: list[str] = []

        # Resolve metadata for the file
        try:
            resp = await client.get(
                f"{_DRIVE_FILES_API}/{file_id}",
                params={"fields": "id,name,mimeType,size,modifiedTime"},
            )
            resp.raise_for_status()
            item = resp.json()
        except Exception as e:
            err = f"Could not fetch metadata for file {file_id}: {e}"
            logger.error("gdrive_single_file_metadata_error", file_id=file_id, error=str(e))
            return ConnectorResult(
                source_id=source_id,
                files=[],
                fetch_metadata={"file_id": file_id},
                errors=[err],
                is_full_sync=True,
            )

        mime = item.get("mimeType", "")
        name_for_gate = str(item.get("name") or "")
        ext_gate = _drive_file_extension(name_for_gate)
        octet_ok = mime in _OCTET_STREAM_MIMES and ext_gate in {
            ".pdf",
            ".docx",
            ".doc",
            ".pptx",
            ".txt",
            ".md",
            ".markdown",
            ".csv",
        }
        if (
            mime not in _EXPORTABLE_GOOGLE_TYPES
            and mime not in _NATIVE_TEXT_TYPES
            and mime not in _OFFICE_TYPES
            and not octet_ok
        ):
            return ConnectorResult(
                source_id=source_id,
                files=[],
                fetch_metadata={"file_id": file_id},
                errors=[f"Unsupported MIME type '{mime}' for file '{item.get('name')}'"],
                is_full_sync=True,
            )

        content, err = await _fetch_file_content(client, item)
        if err:
            errors.append(err)

        files = []
        if content:
            files.append(RawFile(
                path=f"{item['name']} [{file_id}]",
                content=content,
                size_bytes=len(content),
                metadata={
                    "drive_file_id": file_id,
                    "name": item["name"],
                    "modified_time": item.get("modifiedTime"),
                },
            ))

        logger.info("gdrive_single_file_sync_complete", file_id=file_id, fetched=len(files))
        return ConnectorResult(
            source_id=source_id,
            files=files,
            fetch_metadata={
                "file_id": file_id,
                "modified_time": item.get("modifiedTime"),
                "revision_id": (
                    f"drive:{file_id}:{item.get('modifiedTime')}"
                    if item.get("modifiedTime")
                    else None
                ),
            },
            errors=errors,
            is_full_sync=True,
        )

    async def _fetch_full(
        self,
        client: httpx.AsyncClient,
        source_id: str,
        folder_id: str | None,
    ) -> ConnectorResult:
        logger.info("gdrive_full_sync_start", folder_id=folder_id)
        files: list[RawFile] = []
        errors: list[str] = []

        # Build the files list query
        query_parts = [
            "trashed = false",
            "(" + " or ".join(
                f"mimeType = '{m}'"
                for m in list(_EXPORTABLE_GOOGLE_TYPES) + list(_NATIVE_TEXT_TYPES) + list(_OFFICE_TYPES)
            ) + ")",
        ]
        if folder_id:
            query_parts.append(f"'{folder_id}' in parents")

        q = " and ".join(query_parts)
        page_token: str | None = None

        while True:
            params: dict = {
                "q": q,
                "pageSize": 100,
                "fields": "nextPageToken,files(id,name,mimeType,size,modifiedTime)",
            }
            if page_token:
                params["pageToken"] = page_token

            try:
                resp = await client.get(_DRIVE_FILES_API, params=params)
                resp.raise_for_status()
            except Exception as e:
                errors.append(f"Drive file list error: {e}")
                logger.error("gdrive_list_error", error=str(e))
                break

            data = resp.json()
            for item in data.get("files", []):
                content, err = await _fetch_file_content(client, item)
                if err:
                    errors.append(err)
                elif content:
                    files.append(RawFile(
                        path=f"{item['name']} [{item['id']}]",
                        content=content,
                        size_bytes=len(content),
                        metadata={
                            "drive_file_id": item["id"],
                            "name": item["name"],
                            "modified_time": item.get("modifiedTime"),
                        },
                    ))

            page_token = data.get("nextPageToken")
            if not page_token:
                break

        # Get a startPageToken to use for the next delta sync
        next_page_token = await _get_start_page_token(client)

        logger.info(
            "gdrive_full_sync_complete",
            files=len(files), errors=len(errors),
        )
        return ConnectorResult(
            source_id=source_id,
            files=files,
            fetch_metadata={"folder_id": folder_id},
            errors=errors,
            is_full_sync=True,
            next_page_token=next_page_token,
        )

    async def _fetch_delta(
        self,
        client: httpx.AsyncClient,
        source_id: str,
        start_page_token: str,
    ) -> ConnectorResult:
        """Fetch only files changed since the stored pageToken."""
        logger.info("gdrive_delta_sync_start")
        files: list[RawFile] = []
        deleted_paths: list[str] = []
        errors: list[str] = []

        page_token: str | None = start_page_token
        last_token = start_page_token

        while page_token:
            try:
                resp = await client.get(
                    _DRIVE_CHANGES_API,
                    params={
                        "pageToken": page_token,
                        "pageSize": 100,
                        "fields": (
                            "nextPageToken,newStartPageToken,"
                            "changes(removed,fileId,file(id,name,mimeType,size,modifiedTime,trashed))"
                        ),
                    },
                )
                resp.raise_for_status()
            except Exception as e:
                errors.append(f"Drive changes error: {e}")
                logger.error("gdrive_changes_error", error=str(e))
                break

            data = resp.json()

            for change in data.get("changes", []):
                file_info = change.get("file", {})
                file_id = change.get("fileId", "")
                removed = change.get("removed", False) or file_info.get("trashed", False)
                name = file_info.get("name", file_id)
                path = f"{name} [{file_id}]"
                mime = file_info.get("mimeType", "")

                if removed:
                    deleted_paths.append(path)
                    continue

                ext_d = _drive_file_extension(str(name))
                octet_ok_d = mime in _OCTET_STREAM_MIMES and ext_d in {
                    ".pdf",
                    ".docx",
                    ".doc",
                    ".pptx",
                    ".txt",
                    ".md",
                    ".markdown",
                    ".csv",
                }
                if (
                    mime not in _EXPORTABLE_GOOGLE_TYPES
                    and mime not in _NATIVE_TEXT_TYPES
                    and mime not in _OFFICE_TYPES
                    and not octet_ok_d
                ):
                    continue  # skip non-text types

                content, err = await _fetch_file_content(client, file_info)
                if err:
                    errors.append(err)
                elif content:
                    files.append(RawFile(
                        path=path,
                        content=content,
                        size_bytes=len(content),
                        metadata={
                            "drive_file_id": file_id,
                            "name": name,
                            "modified_time": file_info.get("modifiedTime"),
                        },
                    ))

            # Drive returns newStartPageToken on the last page of changes
            new_start = data.get("newStartPageToken")
            if new_start:
                last_token = new_start

            page_token = data.get("nextPageToken")

        logger.info(
            "gdrive_delta_sync_complete",
            changed=len(files), deleted=len(deleted_paths), errors=len(errors),
        )
        return ConnectorResult(
            source_id=source_id,
            files=files,
            fetch_metadata={},
            errors=errors,
            is_full_sync=False,
            deleted_paths=deleted_paths,
            next_page_token=last_token,
        )


# ─── Helpers ──────────────────────────────────────────────────────────────────


async def _get_start_page_token(client: httpx.AsyncClient) -> str | None:
    """Fetch the current Drive Changes startPageToken for use in future delta syncs."""
    try:
        resp = await client.get(
            _DRIVE_START_PAGE_TOKEN_API,
            params={"fields": "startPageToken"},
        )
        resp.raise_for_status()
        return resp.json().get("startPageToken")
    except Exception as e:
        logger.warning("gdrive_start_page_token_error", error=str(e))
        return None


async def _fetch_file_content(
    client: httpx.AsyncClient, item: dict
) -> tuple[str | None, str | None]:
    """
    Fetch or export a Drive file's text content.

    - Google Docs / Sheets / Slides  → export endpoint (text/plain or text/csv)
    - Office formats (DOCX, PPTX)    → download binary, parse locally
    - PDF                            → download binary, pypdf text extraction
    - Plain text / Markdown          → download decoded as UTF-8 text
    """
    file_id = item.get("id", "")
    name = item.get("name", file_id)
    mime = item.get("mimeType", "")

    if not file_id:
        return None, "Missing file ID"

    try:
        if mime in _EXPORTABLE_GOOGLE_TYPES:
            export_mime = _EXPORTABLE_GOOGLE_TYPES[mime]
            resp = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                params={"mimeType": export_mime},
            )
            resp.raise_for_status()
            raw = resp.content
            content = raw.decode("utf-8-sig", errors="replace").replace("\x00", "")
            return content if content.strip() else None, None

        if mime in _OCTET_STREAM_MIMES:
            ext_o = _drive_file_extension(name)
            if ext_o not in {
                ".pdf",
                ".docx",
                ".doc",
                ".pptx",
                ".txt",
                ".md",
                ".markdown",
                ".csv",
            }:
                return None, None

        # Office / PDF / plain / octet-stream (single-file only) — bytes only, never resp.text.
        size = int(item.get("size", 0) or 0)
        if size > MAX_FILE_SIZE_BYTES:
            logger.debug("gdrive_skip_large_file", name=name, mime=mime, size=size)
            return None, None
        resp = await client.get(
            f"https://www.googleapis.com/drive/v3/files/{file_id}",
            params={"alt": "media"},
        )
        resp.raise_for_status()
        data = resp.content

        if mime in _OFFICE_TYPES:
            text = _extract_office_text(data, _OFFICE_TYPES[mime], name)
            return text if text and text.strip() else None, None

        if mime == "application/pdf" or mime in _OCTET_STREAM_MIMES:
            text = _decode_drive_file_body(data, name=name, mime=mime)
            return text if text and text.strip() else None, None

        text = _decode_drive_file_body(data, name=name, mime=mime)
        return text if text and text.strip() else None, None

    except Exception as e:
        logger.warning("gdrive_file_fetch_error", file_id=file_id, name=name, error=str(e))
        return None, f"Failed to fetch '{name}': {e}"


def _extract_office_text(data: bytes, parser: str, name: str) -> str | None:
    """Extract plain text from a downloaded Office binary (DOCX or PPTX)."""
    try:
        if parser == "docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs) or None
        if parser == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            lines: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n\n".join(lines) or None
    except Exception as e:
        logger.warning("gdrive_office_parse_error", name=name, parser=parser, error=str(e))
    return None
